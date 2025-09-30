import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_DIR = Path(r"C:\Users\kened\Desktop\school_xyz\schoolmangmentsystem")
PACKAGE_JSON = BASE_DIR / "package.json"
OUTPUT_NAME_DEFAULT = "Pitch_Deck.pptx"


def load_app_name():
    try:
        data = json.loads(PACKAGE_JSON.read_text(encoding="utf-8"))
        name = data.get("name") or "School Management System"
    except Exception:
        name = "School Management System"
    # Sanitize for filename/display
    display_name = name.strip().title()
    fname = re.sub(r"[^A-Za-z0-9_-]+", "_", display_name)
    return display_name, f"{fname}_Pitch_Deck.pptx"


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bulleted_slide(prs, title, bullets):
    # Title and Content layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    def add_para(text, level=0):
        if len(tf.paragraphs) == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level

    for item in bullets:
        if isinstance(item, str):
            add_para(item, 0)
        elif isinstance(item, (list, tuple)):
            # First element is the parent bullet, rest are sub-bullets
            if not item:
                continue
            parent = item[0]
            add_para(parent, 0)
            for sub in item[1:]:
                add_para(sub, 1)
        elif isinstance(item, dict):
            # {"text": "...", "subs": ["...", "..."]}
            text = item.get("text", "")
            subs = item.get("subs", [])
            add_para(text, 0)
            for sub in subs:
                add_para(sub, 1)


def build_presentation():
    app_name, suggested_output = load_app_name()
    prs = Presentation()

    # Title slide
    add_title_slide(
        prs,
        title=f"{app_name}",
        subtitle="Unified, cross-platform school management platform (Web • Android • iOS)"
    )

    # Problem
    add_bulleted_slide(prs, "Problem", [
        "Schools juggle fragmented tools (attendance, communication, records)",
        "Manual processes lead to errors and delays",
        "Poor data visibility for teachers and admins",
        "Legacy ERPs have clunky UX and slow adoption",
    ])

    # Solution
    add_bulleted_slide(prs, "Solution", [
        f"{app_name}: a modern, modular school OS built with React Native + Expo",
        [
            "Cross-platform: single codebase for Web/Android/iOS",
            "Multi-tenant backend powered by Supabase",
            "Real-time, reliable, and secure",
        ],
        [
            "Teacher-first UX",
            "Fast, responsive lists and forms",
            "Clean, configurable logging for focused debugging",
        ],
    ])

    # Product Highlights
    add_bulleted_slide(prs, "Product Highlights", [
        [
            "Teacher Dashboard",
            "Leave Application: enhanced modal, validation, smooth scrolling",
            "My Students: fast list, search/filter, detail modal with proper scrolling",
        ],
        [
            "Attendance & Records",
            "Optimized flows with clear success/error feedback",
            "Export to PDF (reports) and data utilities",
        ],
        [
            "Data Import",
            "Bulk student import (~700 records) with validation",
            "Default class creation, tenant assignment, schema alignment",
        ],
    ])

    # Market (placeholders – to adjust)
    add_bulleted_slide(prs, "Market", [
        "Target: K-12 schools and private institutions",
        "Pain is universal: attendance, communication, records, reporting",
        "Top of funnel via pilots, local partnerships, and educator communities",
    ])

    # Business Model (placeholders – to adjust)
    add_bulleted_slide(prs, "Business Model", [
        "SaaS, per-student monthly or annual pricing",
        [
            "Tiers: Basic (attendance, student records)",
            "Pro (communications, reporting, PDF exports)",
            "Enterprise (advanced analytics, SSO, priority support)",
        ],
    ])

    # Traction (derived from project artifacts)
    add_bulleted_slide(prs, "Traction", [
        "Core teacher workflows implemented",
        "Clean console via categorized logging across 200+ files",
        "Bulk import tool processed ~700 student records",
        "Foundation for multi-tenant deployment via Supabase",
    ])

    # Go-To-Market
    add_bulleted_slide(prs, "Go-To-Market", [
        [
            "Pilot Programs",
            "Run 2-3 pilots to validate adoption and ROI",
            "Gather testimonials and case studies",
        ],
        [
            "Channel Partners",
            "Local IT resellers and education consultants",
        ],
        [
            "Community & Content",
            "Teacher-led webinars, WhatsApp groups, and demos",
        ],
    ])

    # Competition & Differentiation
    add_bulleted_slide(prs, "Competition & Differentiation", [
        "Legacy ERPs: feature-rich but slow and hard to use",
        "Newer apps: fragmented, mobile-only, or limited",
        [
            "Differentiators",
            "Cross-platform from day one (Expo)",
            "Modern, responsive UX",
            "Real-time Supabase backend, simple tenanting",
            "Configurable logging and maintainable modular code",
        ],
    ])

    # Technology
    add_bulleted_slide(prs, "Technology", [
        [
            "Frontend",
            "React Native (0.81), Expo (SDK 54)",
            "RN Paper, Reanimated, Skia, charts, PDF export",
        ],
        [
            "Backend",
            "Supabase (Auth, DB, Realtime, Storage)",
            "Multi-tenant aware data flows",
        ],
        [
            "DX",
            "Categorized logger, filtered console noise",
            "Optimized lists and forms with platform-aware scrolling",
        ],
    ])

    # Roadmap (placeholders – to adjust)
    add_bulleted_slide(prs, "Roadmap", [
        [
            "Q4",
            "Parent Portal (mobile/web)",
            "Fees & Payments",
            "Messaging & Notifications",
        ],
        [
            "Q1",
            "Exams, Grades, Timetable",
            "Analytics dashboards",
            "Role-based access improvements",
        ],
    ])

    # Team (placeholders)
    add_bulleted_slide(prs, "Team", [
        [
            "Founder/PM",
            "EdTech background, school ops experience",
        ],
        [
            "Engineering",
            "Full-stack RN + Supabase",
            "Mobile & Web performance tuning",
        ],
        [
            "Advisors",
            "School admin, curriculum expert",
        ],
    ])

    # The Ask (placeholders – to adjust)
    add_bulleted_slide(prs, "The Ask", [
        "Seeking pilot partners and early adopters",
        "Optional: raise capital to accelerate product and GTM",
        [
            "Use of funds",
            "Engineering hires",
            "Onboarding & support",
            "GTM experiments",
        ],
    ])

    # Contact
    add_bulleted_slide(prs, "Contact", [
        f"Product: {app_name}",
        "Email: info@example.com",
        "Phone: +1-000-000-0000",
        "Website: https://example.com",
    ])

    # Save
    out_path = BASE_DIR / suggested_output if suggested_output else BASE_DIR / OUTPUT_NAME_DEFAULT
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    path = build_presentation()
    print(f"Pitch deck generated: {path}")
